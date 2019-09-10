# encoding: utf-8

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import requests
import sys
import os

def main():
    
    limpaTela();
    credential = 'ea08bc5e0da9443ad66bdd836ae7ce72';
    artista = '';
    musica = '';
    resposta = '';

    print("####################################################")
    print("####################################################")
    print("##                                                ##")
    print("##  ___                 _        _                ##")
    print("## | _ \___ __ _ __ _  | |   ___| |_ _ _ __ _     ##")
    print("## |  _/ -_) _` / _` | | |__/ -_)  _| '_/ _` |    ##")
    print("## |_| \___\__, \__,_| |____\___|\__|_| \__,_|    ##")
    print("##         |___/                                  ##")                                              
    print("##                                                ##")
    print("##  Busque e salve uma letra de musica em pptx!   ##")
    print("##  Usa API do vagalume                           ##")
    print("##  Author: Lucas Souza                           ##")
    print("##                                                ##")
    print("####################################################")
    print("####################################################")

    artista = input("\nNome do artista/banda: ")
    musica = input("\nNome da musica: ")

  
    #html = urlopen("https://www.letras.mus.br/salvaon/aqui-com-voce/")
    request = requests.get("https://api.vagalume.com.br/search.php?art="+artista+"&mus="+musica+"&apikey="+credential)
    status = request.status_code

    if(status != 200):   
        print("Servidor fora do ar ou dominio incorreto!")
    else:
        #soup = BeautifulSoup(html.read(),"html.parser")
        #for text in soup.find_all("div",{"class":"cnt-letra"}):
            
        #    print(text.get_text('\n'))

        print("\nRealizando pesquisa..")
 
        result = request.json()
       
        global nome_musica
             
        if(result['type'] == 'notfound' or result['type'] == 'song_notfound'):
            print("\nMúsica ou Banda não encontrada..")

            resposta = input("\nDeseja realizar nova pesquisa?(S/N) ")

            if(resposta == 's' or resposta == 'S'):
                limpaTela();
                main();
            else:
                sys.exit();
        else:                         
            print("\nMúsica encontrada!\n")
        
            print(result['mus'][0]['text'])

            resposta = input("\n[?]Essa é a musica desejada?(S/N): ")

            if(resposta == 's' or resposta == 'S'):
                
                nome_musica = result['mus'][0]['name']
                
                print("\n[+] Ok! Gerando apresentação..")

                apresentacao(result['mus'][0]['text']);

                print("[+] Apresentação gerada com sucesso!")

                resposta = input("\nDeseja realizar nova pesquisa?(S/N) ")

                if(resposta == 's' or resposta == 'S'):
                    limpaTela();
                    main();
                else:                                
                    sys.exit();
            else:
                limpaTela();
                main();
            

def apresentacao(texto): #texto , quantidade

    texto = texto.split('\n')
    musica = []

    for linha in texto:
        musica = texto
        
    #num = 0
    #for x in musica:
    #    print(musica[num].strip())
    #    num = num + 1

    y=0
    n_slides = 0
    quantidadePorSlide = 5
    calculoTamanhoTexto = len(musica) / quantidadePorSlide
    restanteTexto = len(musica) % quantidadePorSlide
    total = int(calculoTamanhoTexto)+restanteTexto
    prs = Presentation()
    #print(len(musica))
    for x in range(total):
        
        page_slide = prs.slide_layouts[n_slides]
        slide = prs.slides.add_slide(page_slide)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0,0,0)
        
        top = width = height = Inches(2)
        left = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)

        tf = txBox.text_frame             
        
        
        for i in range(quantidadePorSlide+1):
            if(y < len(musica)):

                p = tf.add_paragraph()
                p.text = musica[y]
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255,255,255)

                                                                   
                #print(musica[y]+'linha:'+str(y))
                y = y+1
            else:
                prs.save(nome_musica+".pptx")
                #x = input('Parou no ELSE aqui')
                return
        print('')
        n_slides = n_slides + 1;

    prs.save(nome_musica+".pptx")
    
def limpaTela():
    os.system('cls' if os.name == 'nt' else 'clear')
    
if __name__ == "__main__":
    main()
